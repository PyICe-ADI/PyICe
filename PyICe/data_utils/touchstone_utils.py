'''
This module is intended to be used in the workflow of analyzing PCB parasitics for power supplies.
The industry standard file type for communicating these parasitics is a touchstone or ".snp" file.

PCB parasitics are found using 3D electro-magnetic solvers, such as Keysight ADS's sub-tool, EM Pro.
The simulation is carried out by selecting ports of a network (positive and negative pins) as physical
locations on a circuit board layout file, attaching z0 terminations to them and finding the S parameter
matrix describing the N-port network over frequency. The output of this 3D simulation is called a touchstone 
or ".snp" file, representing the file extension, where 's' is the parameter type ('s', 'z', 'y', and others),
'n' is replaced by the number of ports, and 'p' means port. This module is meant to be a general purpose
place for functions and workflows that manipulate and analyze touchstones. This workflow is also intended
to pair with ADI's integrated circuit environment, Adice, taking inputs from Adice and doing things such 
as port number reduction which can be more cumbersome or simulation intensive in Adice.

The open source python module, 'scikit-rf' is used heavily in this module for interacting with and
creating touchstones. Scikit-rf is very powerful and can be used to make all sorts of example/reference
touchstones. More useful scripting can be done with scikit-rf than just what is done in this module.

The fundamental operating principle for all of this is a concept used widely in RF engineering,
the "Two-Port Network". A two-port network, or generally an N-port network, is a different way of looking
at circuits as opposed to lumped element models (inductors, resistors, capacitors). Networks are just black
boxes with electrical ports for interaction, where each port has a positive and negative pin, and the
transform from port to port is represented by impedances and dependant sources.This is fundamental to understand 
before using the touchstone utils, and the Wikipedia page for two-port networks is great for this.

A commonly used expression for describing these transformations is S parameters, which come in an N x N x f matrix
where N is the number of ports and f is the number of frequency points. A helpful summary is that S parameters 
are the power gain to a port, from a port. So S21 is the power gain to port 2, from port 1, and S22 is the power 
reflected back to port 2, from port 2.

Author: Ben Leverett
Date Created: 2023-03-23
Version 0.0
'''

import skrf
import numpy             as np
import matplotlib.pyplot as mplt
import scipy.optimize    as optimize
import scipy.signal      as signal
from pathlib             import Path
import datetime
import pptx
import sympy


class touchstone_utils():
    '''
    N-port touchstone. 
    
    The purpose of this class is to simplify and manipulate touchstone files by turning them into
    networks and performing various operations on them. This class uses the open source library "scikit-rf"
    to import touchstone files as well as other operations.
    
    =====================  =============================================
    Property               Meaning
    =====================  =============================================
    :attr:`network`        current network object for operations
    :attr:`z0`             characteristic impedance of the touchstone (assuming uniform z0)
    :attr:`freqs`          frequency vector, numpy array
    =====================  =============================================
    '''
    def __init__(self, file_name=None):
        '''
            Creates a touchstone utils object.
        Takes in an N port touchstone file for manipulation and plotting.
        One of the main goals of the touchstone utils is to simplify networks from N port to M port.
        Touchstone utils can also plot bode plots and create networks and touchstones files.
        
        Parameters
        -----------
        file_name: str
            Touchstone file to be imported, extension ".s?p"
        '''
        self.network = None

        if file_name != None:  #TODO: More error checking
            self.network = skrf.Network(file_name)

    def __del__(self):
        '''
         Deconstructor. Called when the object is destroyed, usually when use script terminates.
         TODO: this did not work, maybe there is a better way to call mplt.show() when exiting scripts
        '''
        mplt.show()

    def output_touchstone(self, output_file_name, output_file_dir, input_network=None, start=None, stop=None, output_form='ri'):    
        '''
        Outputs a touchstone of a specified network. If none specified, the current touchstone_utils network is used.
        
        Parameters
        -----------
        output_file_name: str
            Output file name without file extension. File extension will be assigned.
        output_file_dir: str
            Output file directory.
        input_network: skrf.network.Network, optional
            Optional input network to output a touchstone of, if none given, the current touchstone utils network is used
        start: float, optional
            Start frequency for the data of the network. If none given or lower than output_network.f[0] it is ignored.
        stop: float, optional
            Stop frequency for the data of the network. If none given or higher than network.f[-1] it is ignored.
        output_form: str, optional
            Output form based on skrf.network.Network.write_touchstone.
            Common forms are 'ri' for real & imaginary and 'ma' for magnitude & angle (degrees).
        '''
        # TODO: Add model type option to choose S, Y, or Z parameters.
        # TODO: Make the output file name automatic based on the current or input network
        # TODO: Look at touchstone from the LTC3315 customer and see if their touchstone form phase is in degrees or radians
        # TODO: Comment the touchstone with some ADI copyright stuff
        output_network = input_network
        if output_network == None:
            output_network = self.network.copy()
        if start != None:
            output_network.crop(start, output_network.f[-1])
        if stop != None:
            output_network.crop(output_network.f[0], stop)        
        output_network.write_touchstone(filename = output_file_name, 
                                        dir = output_file_dir, 
                                        skrf_comment = False,
                                        form = output_form, # Real Imaginary, could also use mag for Magnitude Phase(deg)
                                       )

    def network_Nport_to_Mport(self, port_nums):
        '''
        Converts current touchstone utils network to an M-port network of ports 'port_nums'.
        
        Parameters
        -----------
        port_nums: list(type(int))
            Port numbers from the current N-port network to be used for M-port network.
            M-port port numbers 1-M are assigned in order of port_nums, so port_nums = [2, 1] will give
                M-port port 1 as N-port port 2 and M-port port 2 as N-port port 1.
            Common forms are 'ri' for real & imaginary and 'ma' for magnitude & angle (degrees).
        '''
        # TODO: Add input network
        # TODO: Add error handling for port_nums outside of N-port port count
        def open_terminator(name): #implicit freqs, z0
            open_port = skrf.Circuit.Port(frequency=freqs, name='open_term_port', z0=z0)
            open_series = skrf.Circuit.Open(frequency=freqs, name='open', z0=z0)
            shunt =  skrf.Circuit.Ground(frequency=freqs, name='GND', z0=z0)
            open_term_cnx = [[(open_port, 0), (open_series, 0)],
                             [(open_series, 1), (shunt, 0)],
                            ]
            cir_open_term = skrf.Circuit(open_term_cnx)
            nw_open_term = cir_open_term.network
            nw_open_term.name = name
            return nw_open_term
        n = self.network.nports
        z0 = self.network.z0[0][0]
        freqs = self.network.frequency
        capped_connections = []
        unused_port_nums = list(range(1, n+1))
        port_name_counter = 0
        for port_num in port_nums:
            capped_connections.append([(self.network, port_num-1), (skrf.Circuit.Port(frequency=freqs, name=f"P{port_name_counter+1}", z0=z0), 0)])
            unused_port_nums.remove(port_num)
            port_name_counter += 1
        for unused_port in unused_port_nums:
            capped_connections.append([(self.network, unused_port-1), (open_terminator(f"open_term{unused_port}"), 0)])
        capped_circuit = skrf.Circuit(capped_connections)
        capped_circuit.plot_graph(network_labels=True,
               port_labels=True,
               edge_labels=True,
              )
        output_network = skrf.network.Network(frequency=freqs, s=capped_circuit.network.s, z0=z0, name = "M Port Network")
        return output_network # TODO: See about using the Network constructor the way it is used in test_make_Nport_Mport_reference_networks
        
    def network_Nport_to_1port(self, source_port_num, load_port_num):
        '''
        Converts current touchstone utils network to an 1-port network.
        The input port is the source_port_num and which port port is shorted to make the 1-port network is load_port_num.
        
        Parameters
        -----------
        source_port_num: int
            Port number from the current N-port network to be used as the input port for the 1-port network.
        output_port_num: int
            Port number from the current N-port network to be shorted for the simplification to a 1-port network.
        '''
        # TODO: Add functionality for load_port_num to be multiple load ports that get shorted.
        def short_terminator(name): #implicit freqs
            short_port = skrf.Circuit.Port(frequency=freqs, name='short_term_port', z0=z0)
            shunt =  skrf.Circuit.Ground(frequency=freqs, name='GND', z0=z0)
            short_term_cnx = [[(short_port, 0), (shunt, 0)],
                             ]
            cir_short_term = skrf.Circuit(short_term_cnx)
            nw_short_term = cir_short_term.network
            nw_short_term.name = name
            return nw_short_term
        z0 = self.network.z0[0][0]
        freqs = self.network.frequency
        # TODO: Have to do the trick with the Network constructor and getting the s data to get around error
        reduced_2port = self.network_Nport_to_Mport([source_port_num, load_port_num])
        new_source_port = skrf.Circuit.Port(frequency=freqs, name='source', z0=z0)
        ground =        skrf.Circuit.Ground(frequency=freqs, name='gnd',    z0=z0)
        connections_reduce_2_to_1_port = [[(reduced_2port, 0), (new_source_port, 0)],
                                          [(reduced_2port, 1), (ground, 0)],
                                         ]
        circuit_1port = skrf.Circuit(connections_reduce_2_to_1_port)
        circuit_1port.plot_graph(network_labels=True,
                                  port_labels=True,
                                  edge_labels=True,
                                 )
        output_network = skrf.network.Network(frequency=freqs, s=circuit_1port.network.s, z0=z0, name = "Test 1 Port Network")
        return output_network
        
        
    def set_network(self, network):
        '''
        TODO: Make an error checking network setter.
        '''
        pass
        
    def test_make_Nport_Mport_reference_networks(self, N, M=None):
        '''
        Converts current touchstone utils network to an 1-port network.
        The input port is the source_port_num and which port port is shorted to make the 1-port network is load_port_num.
        
        Parameters
        -----------
        source_port_num: int
            Port number from the current N-port network to be used as the input port for the 1-port network.
        output_port_num: int
            Port number from the current N-port network to be shorted for the simplification to a 1-port network.
        '''
        # TODO: Complete the M-port portion of this function. Currently just made N-port and used Adice to verify the simplification.
        # TODO: Find a way to compare the networks locally without Adice
        
        # Make N port and M port touchstones of a circuit normally
        # Output the two touchstones and diff
        # Use python 'difflib' to diff the two touchstone files
        # Or use Adice or some other simulators to compare the networks
        def port_terminator(name):
            return skrf.Circuit.Port(frequency=freqs, name=name,  z0=z0)
        def series_resistor(name, resistance):
            return skrf.Circuit.SeriesImpedance(frequency=freqs, name=name,  z0=z0, Z=resistance)
        freqs = skrf.frequency.Frequency(start=10, stop=10e6, 
                                         npoints=1000, unit='Hz', 
                                         sweep_type='log'
                                        )
        z0 = 0.1
        connections_Nport = []
        resistors = []
        middle_resistor_node = []
        for i in range(1, N + 1): 
            resistors.append(series_resistor(name = f"Res{i}", resistance = i * z0))
            connections_Nport.append([(port_terminator(name = f"Port{i}"), 0), (resistors[i-1], 0)])
            middle_resistor_node.append((resistors[i-1], 1))
        connections_Nport.append(middle_resistor_node)
        circuit = skrf.Circuit(connections_Nport)
        circuit.plot_graph(network_labels=True,
                           port_labels=True,
                           edge_labels=True,
                          )
        output_network = skrf.network.Network(frequency=freqs, s=circuit.network.s, z0=z0, name = "Test N Port Network")
        return output_network

def sweep_plots_to_pptx(ts_plots_dir, output_pptx_path, date_time_flag=1):
    '''
    Intended to collect plots from an archive folder and automatically add them to a Microsoft
    PowerPoint presentation.
    
    Searches through a directory for PNG or SVG files, creates a PowerPoint and adds them to it.
    The titles of the slides is the filename of the image. There is no additional slides in the PowerPoint.
    The slides have no background.
    
    Parameters
    -----------
        ts_plots_dir: string
            Directory of target touchstone plots to be swept into PowerPoint presentation.
        output_pptx_path: string
            Full file path of output pptx file. The directory must exist, but the ouptut file can be created.
        date_time_flag: bool or int
            Flag for including the current date and time at the beginning of the ouptut pptx file name.
        
    Future Work
    ------------
        TODO: Error checking and feedback
        TODO: Make it look for SVG as well or just SVG
        TODO: Add a way to input presentation templates
        TODO Ben: Change adice script to print SVG
        TODO: Make a from archive workflow that can take the datetime the plots were recorded not the current datetime.
    '''
    
    # TODO: Check that ts_plots_dir exists and has an image
    # TODO: Check that output_pptx_path.parent exists, its okay for the file not to exist.
    # TODO: Check the fomat of date_time_flag

    pptx_output = pptx.Presentation()
    pptx_output.slide_height = pptx.util.Inches(7.5)
    pptx_output.slide_width  = pptx.util.Inches(13.333)
    
    plot_paths_list = Path(arctic_folder + '23_07_06__11_25_35_ltc3315_sparam_report/').glob('**/*.png')
    for plot_path in plot_paths_list:
        plot_file_name = str(plot_path)

        blank_slide_layout = pptx_output.slide_layouts[6]
        slide = pptx_output.slides.add_slide(blank_slide_layout)

        left = pptx.util.Inches(0.36)
        top = pptx.util.Inches(0.18)
        width = pptx.util.Inches(8)
        height = pptx.util.Inches(1)
        title_textbox = slide.shapes.add_textbox(left, top, width, height)
        title_textframe = title_textbox.text_frame
        title_paragraph = title_textframe.paragraphs[0]
        title_paragraph.font.size = pptx.util.Pt(44)
        title_paragraph.font.name = "Calibri Light"
        title_paragraph.text = plot_path.stem

        position_top = pptx.util.Inches(1.3)
        position_left = pptx.util.Inches(0)
        plot_width = pptx.util.Inches(11.16)
        plot_height = pptx.util.Inches(6.2)
        pic = slide.shapes.add_picture(plot_file_name, position_left, position_top, width=plot_width, height=plot_height)
    
    curr_date_time = datetime.datetime.now()
    date_time_string = curr_date_time.strftime("%y_%m_%d__%H_%M_%S_")
    pptx_output.save(arctic_folder + f'/pptx_outputs/{date_time_string}python_pptx_test.pptx')
    
def resistor_ladder_coefficient(r_dc, r_hf, num_stages):
    r_coeff = sympy.symbols('a')
    r_inv = 1 / r_hf 
    for i in range(num_stages):
        r_inv = r_inv + 1 / (r_hf * r_coeff ** (i + 1))
    r_inv = sympy.simplify(r_inv)
    r_eq = sympy.simplify(1 / r_inv)
    equation = r_eq - r_dc
    r_coeff_val = sympy.nsolve(equation, 1)
    return r_coeff_val
    
def inductor_ladder_coefficient(z_hf, z_lf, r_hf, r_coeff, num_stages):
    pass
    
    
def _parallel(Ra, Rb):
    inv_par = sympy.simplify(1/Ra + 1/Rb)
    return sympy.simplify(1/inv_par)
    # return sympy.simplify(sympy.simplify(Ra * Rb) / sympy.simplify(Ra + Rb))
    
###################################################################
###################################################################
###################################################################
###################################################################
################## End of Supported Development ###################
###################################################################
###################################################################
###################################################################
###################################################################



    def dev_make_series_LR_model(self, outputFileName=None, R=1e3, L=1e-3):
        #Makes a zero for Z(1,1) at 160K Hz
        freqs = skrf.frequency.Frequency(start=10, stop=10e6, 
                                         npoints=1000, unit='Hz', 
                                         sweep_type='log'
                                        )
        source_port = skrf.Circuit.Port(frequency=freqs,    name='in',  z0=50)
        # load_port = skrf.Circuit.Port(frequency=freqs,      name='out', z0=50)
        res = skrf.Circuit.SeriesImpedance(frequency=freqs, name='res', z0=50, Z=R)
        ind = skrf.Circuit.SeriesImpedance(frequency=freqs, name='ind', z0=50, Z=1j*freqs.w*L)
        ground =  skrf.Circuit.Ground(frequency=freqs,      name='gnd', z0=50)
        connections = [[(source_port, 0), (res, 0)],
                       [(res, 1),         (ind, 0),],
                       [(ind, 1),         (ground, 0)]
                      ]
        circuit = skrf.Circuit(connections)
        network = circuit.network
        self.network = network
        # z = skrf.network.s2z(network.s, z0=50)
        # print(z[:,0,0])  #Printing Z(1,1) over frequency
        circuit.plot_graph(network_labels=True, #Plot network graph
                           port_labels=True,
                           edge_labels=True,
                          )
        mplt.figure(2)
        network.plot_z_mag(m=0, n=0, logx=True) #Plot Z(1,1) magnitude over frequency
        
    def dev_series_LR_curve_fit(self):
        # if self.network == None:
            # print('No network to fit\n')
            # return
        z11 = np.abs(self.network.z[:,0,0])
        freqs = self.network.f
        def estimator_fun(x, R, L):
            return np.abs(R + 1j*2*3.14*x*L)
        params, covariance = optimize.curve_fit(f=estimator_fun,
                                                xdata=freqs,
                                                ydata=z11,
                                               )
        print(f"R:{params[0]}, L:{params[1]}")
        return params
        
    def dev_curve_fit_3stage_ladder(self, points):
        range = int(points/2)
        z11 = self.network.z[:,0,0]
        freqs = self.network.f
        zero_array = np.zeros(len(freqs))
        index_f0 = np.where(np.angle(z11,deg=True) >= 45)[0][0]
        z11_f0_sample = z11[(index_f0-range):(index_f0+range)]
        freqs_f0_sample = freqs[(index_f0-range):(index_f0+range)]
        zero_array_sample = zero_array[0:points]
        # mplt.figure()
        # mplt.plot(freqs_f0_sample, np.abs(z11_f0_sample), label='ts sample')
        # mplt.yscale('log')
        # mplt.xscale('log')
        # mplt.ylabel('Impedance Magnitude (log|Ohms|)')
        # mplt.xlabel('Frequency (log|Hz|)')
        # mplt.legend()
        
        # mplt.figure()
        # mplt.plot(freqs_f0_sample, np.angle(z11_f0_sample,   deg=True),  label='ts sample')
        # mplt.xscale('log')
        # mplt.ylabel('Impedance Phase (Degrees)')
        # mplt.xlabel('Frequency (log|Hz|)')
        # mplt.legend()
        print(f"index_f0:{index_f0}, f0: {freqs[index_f0]}")
        def parallel(z1,z2):
            return z1*z2/(z1+z2)
        def estimator_fun(x, L0, R1, L1, LL, RR):
            s = x*2j*3.14
            ladder = s*L1*(LL**2) + R1*(RR**3)
            ladder = parallel(ladder, R1*(RR**2))
            ladder = ladder + s*L1*LL
            ladder = parallel(ladder, R1*RR)
            ladder = ladder + s*L1
            ladder = parallel(ladder, R1)
            series_impedance = s*L0 + ladder
            return np.abs(series_impedance-z11_f0_sample)
        params1, covariance1 = optimize.curve_fit(f=estimator_fun,
                                                xdata=freqs_f0_sample,
                                                ydata=zero_array_sample,
                                                bounds=([0,0,0,0,0],[np.inf,np.inf,np.inf,np.inf,np.inf]),
                                               )
        print(f"L0:{params1[0]}, R1:{params1[1]}, L1:{params1[2]}, LL:{params1[3]}, RR:{params1[4]}")
        return(params1)

    def dev_curve_fit_2stage_ladder(self, fmax):
        z11 = self.network.z[:,0,0]
        freqs = self.network.f
        zero_array = np.zeros(len(freqs))
        
        index_fmax = np.where(freqs >= fmax)[0][0] - 1
        z11_thru_fmax = z11[:(index_fmax)]
        freqs_thru_fmax = freqs[:(index_fmax)]
        zero_array_sample = zero_array[0:index_fmax]
        
        weight = 1-np.log10(freqs-1e3+1)/np.log10(freqs[len(freqs)-1])
        # size = len(freqs)
        # temp = 0
        # for i in range(len(freqs)):
            # temp += 1/size
            # weight[size-1-i] = temp
        # mplt.figure()
        # mplt.plot(freqs, weight,  label='weight')
        # mplt.xscale('log')
        # mplt.ylabel('Weight')
        # mplt.xlabel('Frequency (Hz)')
        # mplt.legend()
        
        def parallel(z1,z2):
            return z1*z2/(z1+z2)
        def estimator_fun(x, L0, R1, L1, R2, L2, R3):
            s = x*2j*3.14
            ladder = s*L1 + R2
            ladder = parallel(ladder, R1)
            series_impedance = s*L0 + ladder
            return np.abs(series_impedance-z11) * weight
        params3, covariance3 = optimize.curve_fit(f=estimator_fun,
                                        xdata=freqs,
                                        ydata=zero_array,
                                        p0=[3e-9,35e-3,30e-9,20e-3,100e-9,10e-3],
                                        bounds=([0,0,0,0,0,0],[np.inf,np.inf,np.inf,np.inf,np.inf,np.inf]),
                                       )
        print(f"L0:{params3[0]}, R1:{params3[1]}, L1:{params3[2]}, R2:{params3[3]}, L2:{params3[4]}, R3:{params3[5]}")
        return(params3)     
        
    def dev_curve_fit_1stage_ladder(self, points):
        range = int(points/2)
        z11 = self.network.z[:,0,0]
        freqs = self.network.f
        zero_array = np.zeros(len(freqs))
        index_f0 = np.where(np.angle(z11,deg=True) >= 45)[0][0]
        z11_f0_sample = z11[(index_f0-range):(index_f0+range)]
        freqs_f0_sample = freqs[(index_f0-range):(index_f0+range)]
        zero_array_sample = zero_array[0:points]
        print(f"index_f0:{index_f0}, f0: {freqs[index_f0]}")
        def parallel(z1,z2):
            return z1*z2/(z1+z2)
        def estimator_fun(x, L0, R1, L1, R2):
            s = x*2j*3.14
            ladder = s*L1 + R2
            ladder = parallel(ladder, R1)
            series_impedance = s*L0 + ladder
            return np.abs(series_impedance-z11_f0_sample)
        params2, covariance2 = optimize.curve_fit(f=estimator_fun,
                                                xdata=freqs_f0_sample,
                                                ydata=zero_array_sample,
                                                bounds=([0,0,0,0],[np.inf,np.inf,np.inf,np.inf]),
                                               )
        print(f"L0:{params2[0]}, R1:{params2[1]}, L1:{params2[2]}, R2:{params2[3]}")
        return(params2)
        
    def dev_minimize_2stage_ladder(self, fmax):
        z11 = self.network.z[:,0,0]
        # freqs = self.network.f
        # zero_array = np.zeros(len(freqs))
        
        # index_fmax = np.where(freqs >= fmax)[0][0] - 1
        # z11_thru_fmax = z11[:(index_fmax)]
        # freqs_thru_fmax = freqs[:(index_fmax)]
        # zero_array_sample = zero_array[0:index_fmax]
                
        # def parallel(z1,z2):
            # return z1*z2/(z1+z2)
        # def estimator_fun(variables):
            # L0, R1, L1, R2, L2, R3 = variables
            # s = freqs_thru_fmax*2j*3.14
            # ladder = s*L1 + R2
            # ladder = parallel(ladder, R1)
            # series_impedance = s*L0 + ladder
            # return (np.abs(series_impedance-z11_thru_fmax)**2).sum()
            
        # x, success = optimize.minimize(
    
        
        # params3, covariance3 = optimize.curve_fit(f=estimator_fun,
                                        # xdata=freqs_thru_fmax,
                                        # ydata=zero_array_sample,
                                        # bounds=([0,0,0,0,0,0],[np.inf,np.inf,np.inf,np.inf,np.inf,np.inf]),
                                       # )

        # return x, success
        
      
    def dev_plot_parallel(Req):
        def dev_parallel_solver(Req, R2):
            R1 = R2*Req/(R2-Req)        
            return R1        
        R1 = np.linspace(12e-3,50e-3,100)
        R2 = np.zeros(len(R1))
        for i in range(len(R1)):
            R2[i] = dev_parallel_solver(Req, R1[i])
        print(R1)
        print(R2)
        mplt.figure()
        mplt.plot(R1, R2, label='R1||R2=10m Ohm')
        mplt.ylabel('R2')
        mplt.xlabel('R1')
        mplt.legend()
        
    def dev_plot_log_spaces():
        freq = np.logspace(3, 7, 41)
        y = freq
        # y[10000] = 1
        # y[13300] = 1
        # y[17780] = 1
        # y[23700] = 1
        print(freq)
        mplt.figure()
        mplt.plot(freq, y, label='R1||R2=10m Ohm')
        mplt.ylabel('y')
        mplt.xlabel('freq (log|Hz|)')
        mplt.xscale('log')
        mplt.yscale('log')
        mplt.legend()
        
    def dev_partial_fractions():
        # Make a transfer function that we know has a real pole and a complex pole pair
        # See if we get complex coefficients. There may be another unexpected solution
        # Alternatively, just try to input complex coefficients
        b = [0,1,2]
        a = [3,4,5]
        print(signal.residue(b, a))

    def dev_solve(): #Not used, imported SymPy for this but haven't needed it.
        R1 = .045
        L0 = 7.16e-9
        p2 = 3.16e5
        z2 = 1e5
        p1 = 3.16e4
        z1 = 1e4
        R2, R3, L1, L2 = symbols('R2,R3,L1,L2')
        eq1 = (R1+R2)/(2*np.pi*L1)-p2
        eq2 = R2/(2*np.pi*L1)-z2
        eq3 = (R2+R3)/(2*np.pi*L2)-p1
        eq4 = R3/(2*np.pi*L2)-z1
        eq5 = (1/((1/R1)+(1/R2)+(1/R3)))-.01
        res = solve((eq1,eq2,eq3,eq4,eq5))
        print(res)
        
    def dev_plot_series_LR_error(self, params):
        R, L = params
        freq_obj = self.network.frequency
        freqs = self.network.f
        in_port = skrf.Circuit.Port(       frequency=freq_obj, name='in',  z0=50)
        ind = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L',  z0=50, Z=1j*freq_obj.w*L)
        res = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R',  z0=50, Z=R)
        ground =  skrf.Circuit.Ground(     frequency=freq_obj, name='gnd', z0=50)
        connections = [[(in_port, 0), (ind, 0)],
                       [(ind, 1),     (res, 0)],
                       [(res, 1),     (ground, 0)],
                      ]
        series_LR_circuit = skrf.Circuit(connections)
        series_LR_network = series_LR_circuit.network
        # series_LR_circuit.plot_graph(network_labels=True,
                           # port_labels=True,
                           # edge_labels=True,
                          # )
        network_z11 = self.network.z[:,0,0]
        series_LR_z11 = series_LR_network.z[:,0,0]
        mplt.figure()
        mplt.plot(freqs, np.abs(network_z11),    label='touchstone')
        mplt.plot(freqs, np.abs(series_LR_z11),  label='series_LR')
        mplt.yscale('log')
        mplt.xscale('log')
        mplt.ylabel('Impedance Magnitude (log|Ohms|)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()
        
        mplt.figure()
        mplt.plot(freqs, np.angle(network_z11,   deg=True),  label='touchstone')
        mplt.plot(freqs, np.angle(series_LR_z11, deg=True),  label='series_LR')
        mplt.xscale('log')
        mplt.ylabel('Impedance Phase (Degrees)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()
    
    def dev_plot_3stage_ladder(self, params):
        L0, R1, L1, LL, RR = params
        freq_obj = self.network.frequency
        freqs = self.network.f
        in_port = skrf.Circuit.Port(frequency=freq_obj,         name='in',  z0=50)
        ind0 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L0',  z0=50, Z=1j*freq_obj.w*L0)
        res1 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R1',  z0=50, Z=R1)
        ind1 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L1',  z0=50, Z=1j*freq_obj.w*L1)
        res2 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R2',  z0=50, Z=R1*RR)
        ind2 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L2',  z0=50, Z=1j*freq_obj.w*L1*LL)
        res3 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R3',  z0=50, Z=R1*RR**2)
        ind3 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L3',  z0=50, Z=1j*freq_obj.w*L1*LL**2)
        res4 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R4',  z0=50, Z=R1*RR**3)
        ground =  skrf.Circuit.Ground(      frequency=freq_obj, name='gnd', z0=50)
        connections = [[(in_port, 0), (ind0, 0)],
                       [(ind0, 1), (res1, 0), (ind1, 0)],
                       [(ind1, 1), (res2, 0), (ind2, 0)],
                       [(ind2, 1), (res3, 0), (ind3, 0)],
                       [(ind3, 1), (res4, 0)],
                       [(res1, 1), (res2, 1), (res3, 1), (res4, 1), (ground, 0)],
                      ]
        ladder_circuit = skrf.Circuit(connections)
        ladder_network = ladder_circuit.network
        # ladder_circuit.plot_graph(network_labels=True,
                           # port_labels=True,
                           # edge_labels=True,
                          # )
        
        network_z11 = self.network.z[:,0,0]
        ladder_z11 = ladder_network.z[:,0,0]
        mplt.figure()
        mplt.plot(freqs, np.abs(network_z11), label='touchstone')
        mplt.plot(freqs, np.abs(ladder_z11),  label='3stage_ladder')
        mplt.yscale('log')
        mplt.xscale('log')
        mplt.ylabel('Impedance Magnitude (log|Ohms|)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()
        
        mplt.figure()
        mplt.plot(freqs, np.angle(network_z11, deg=True), label='touchstone')
        mplt.plot(freqs, np.angle(ladder_z11, deg=True),  label='3stage_ladder')
        mplt.xscale('log')
        mplt.ylabel('Impedance Phase (Degrees)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()
        
    def dev_plot_2stage_ladder(self, params):
        L0, R1, L1, R2, L2, R3 = params
        freq_obj = self.network.frequency
        freqs = self.network.f
        in_port = skrf.Circuit.Port(frequency=freq_obj,         name='in',  z0=50)
        ind0 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L0',  z0=50, Z=1j*freq_obj.w*L0)
        res1 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R1',  z0=50, Z=R1)
        ind1 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L1',  z0=50, Z=1j*freq_obj.w*L1)
        res2 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R2',  z0=50, Z=R2)
        ind2 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L2',  z0=50, Z=1j*freq_obj.w*L2)
        res3 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R3',  z0=50, Z=R3)
        ground =  skrf.Circuit.Ground(      frequency=freq_obj, name='gnd', z0=50)
        connections = [[(in_port, 0), (ind0, 0)],
                       [(ind0, 1), (res1, 0), (ind1, 0)],
                       [(ind1, 1), (res2, 0), (ind2, 0)],
                       [(ind2, 1), (res3, 0)],
                       [(res1, 1), (res2, 1), (res3, 1), (ground, 0)],
                      ]
        ladder_circuit = skrf.Circuit(connections)
        ladder_network = ladder_circuit.network
        # ladder_circuit.plot_graph(network_labels=True,
                           # port_labels=True,
                           # edge_labels=True,
                          # )
        network_z11 = self.network.z[:,0,0]
        ladder_z11 = ladder_network.z[:,0,0]
        
        mplt.figure()
        mplt.plot(freqs, np.angle(network_z11, deg=True), label='touchstone')
        mplt.plot(freqs, np.angle(ladder_z11, deg=True),  label='2stage_ladder')
        mplt.xscale('log')
        mplt.ylabel('Impedance Phase (Degrees)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()
        
        mplt.figure()
        mplt.plot(freqs, np.abs(network_z11), label='touchstone')
        mplt.plot(freqs, np.abs(ladder_z11),  label='2stage_ladder')
        mplt.yscale('log')
        mplt.xscale('log')
        mplt.ylabel('Impedance Magnitude (log|Ohms|)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()
        
    def dev_plot_1stage_ladder(self, params):
        L0, R1, L1, R2 = params
        freq_obj = self.network.frequency
        freqs = self.network.f
        in_port = skrf.Circuit.Port(frequency=freq_obj,         name='in',  z0=50)
        ind0 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L0',  z0=50, Z=1j*freq_obj.w*L0)
        res1 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R1',  z0=50, Z=R1)
        ind1 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='L1',  z0=50, Z=1j*freq_obj.w*L1)
        res2 = skrf.Circuit.SeriesImpedance(frequency=freq_obj, name='R2',  z0=50, Z=R2)
        ground =  skrf.Circuit.Ground(      frequency=freq_obj, name='gnd', z0=50)
        connections = [[(in_port, 0), (ind0, 0)],
                       [(ind0, 1), (res1, 0), (ind1, 0)],
                       [(ind1, 1), (res2, 0)],
                       [(res1, 1), (res2, 1), (ground, 0)],
                      ]
        ladder_circuit = skrf.Circuit(connections)
        ladder_network = ladder_circuit.network
        # ladder_circuit.plot_graph(network_labels=True,
                           # port_labels=True,
                           # edge_labels=True,
                          # )
        
        network_z11 = self.network.z[:,0,0]
        ladder_z11 = ladder_network.z[:,0,0]
        mplt.figure()
        mplt.plot(freqs, np.abs(network_z11), label='touchstone')
        mplt.plot(freqs, np.abs(ladder_z11),  label='1stage_ladder')
        mplt.yscale('log')
        mplt.xscale('log')
        mplt.ylabel('Impedance Magnitude (log|Ohms|)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()
        
        mplt.figure()
        mplt.plot(freqs, np.angle(network_z11, deg=True), label='touchstone')
        mplt.plot(freqs, np.angle(ladder_z11, deg=True),  label='1stage_ladder')
        mplt.xscale('log')
        mplt.ylabel('Impedance Phase (Degrees)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()

    
    def dev_LR_minimize_fit(self):
        if self.network == None:
            print('No network to fit\n')
            return
        z11 = self.network.z[:,0,0]
        freqs = self.network.f
        def error_fun(params):
            print(params.shape)
            R, L = params
            tf = R + 2j*3.14*freqs*L
            return (tf-z11)**2

        # res_guess = np.linspace(start=1e2,  stop=1e5,  num=100)
        # ind_guess = np.linspace(start=1e-2, stop=1e-5, num=100)
        # variables = np.array((res_guess, ind_guess)).T
        # print(variables.shape)
        guess = [150, 1.5e-3]
        optimize_result = optimize.minimize(fun=error_fun, x0=guess, bounds=[(1e2,1e5),(1e-5,1e-2)])
        print(optimize_result) 
        
        # print(variables)
        # print(variables[:,0])
        # print(variables[:,1])
        # print(self.network.z.shape)
        #### 
    
    def dev_aprx_ckt_values(corners, Rdc, Rhf):
        z1, p1, z2, p2, z3 = corners
        def parallel(z1,z2):
            return z1*z2/(z1+z2)
        R1 = Rhf
        L0 = R1/(2*3.14*z3)
        R2 = R1/(p2/z2-1)
        L1 = R2/(2*3.14*z2)
        R3 = parallel(R1,R2)*Rdc/(parallel(R1,R2)-Rdc)
        L2 = (R2+R3)/(2*3.14*p1)
        #z1 not used
        print(f"R1={R1}; R2={R2}; R3={R3}; L0={L0}; L1={L1}; L2={L2}")
        return [L0,R1,L1,R2,L2,R3]

    def dev_S11_coefficients_Z11(y, x):
        #Z11 = Zo*(1+S11)/(1-S11)
        #TODO: instead of returning the numerator and denominator coefficients, just use the substitution stuff
        #   to put in values for s and plot locally
        s = sympy.symbols('s')
        z0 = 0.1
        s11_numerator = 0
        s11_denominator = 0
        for i in range(len(y)):
            s11_numerator += y[i]*(s**(len(y)-i-1))
        for j in range(len(x)):
            s11_denominator += x[j]*(s**(len(x)-j-1)) 
        print(f"s11_num:{s11_numerator}") 
        print(f"s11_den:{s11_denominator}") 
        # s11_numerator = y[0]*(s**4)+y[1]*(s**3)+y[2]*(s**2)+y[3]*s+y[4]
        # s11_denominator = x[0]*(s**4)+x[1]*(s**3)+x[2]*(s**2)+x[3]*s+x[4]
        # print(f"s11_num:{s11_numerator}") 
        # print(f"s11_den:{s11_denominator}") 
        
        # divi, rem = sympy.div(numerator, denominator, domain='QQ')
        # print(f"S11 = {rem.evalf(5)}")  #+divi.evalf()

        s11 = s11_numerator/s11_denominator
        
        top = sympy.simplify(1-s11)
        # top = top*z0
        top = sympy.simplify(top)
        bot = sympy.simplify(z0*(1+s11))
        
        z11 = bot/top
        print(f"Z11 = {z11}\n")
        
        z_num, z_den = sympy.fraction(z11)
        print(f"numerator = {z_num}")
        print(f"denominator = {z_den}")
        
        z_num_poly = sympy.Poly(z_num)
        z_num_coeff = z_num_poly.all_coeffs()
        # print(f"num_coeff = {z_num_coeff}")
        
        z_den_poly = sympy.Poly(z_den)
        z_den_coeff = z_den_poly.all_coeffs()
        # print(f"den_coeff = {z_den_coeff}")
        
        r, p, k = signal.residue(z_num_coeff, z_den_coeff)
        
        return z_num_coeff, z_den_coeff, r, p, k

    def dev_plot_rational_coefficients(self, y, x):
        freqs = np.logspace(3, 8, num=51)
        s = 2j*3.14*freqs
        
        # num = [y[0]*(s**4)+y[1]*(s**3)+y[2]*(s**2)+y[3]*s+y[4]]
        # den = [x[0]*(s**4)+x[1]*(s**3)+x[2]*(s**2)+x[3]*s+x[4]]
        for i in range(len(y)):
            m11_numerator += y[i]*(s**(len(y)-i-1))
        for j in range(len(x)):
            m11_denominator += x[j]*(s**(len(x)-j-1))
        Z11 = np.divide(num, den)
    
        freqs2 = np.reshape(freqs, (1,51))

        mplt.figure()
        mplt.plot(self.network.f, np.abs(self.network.z[:,0,0]), label='touchstone')
        mplt.plot(np.transpose(freqs2), np.transpose(np.abs(Z11)), label='curve_fit')
        mplt.yscale('log')
        mplt.xscale('log')
        mplt.ylabel('Impedance Magnitude (log|Ohms|)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()
        
        # mplt.figure()
        # mplt.plot(self.network.f, np.angle(self.network.z[:,0,0]), label='touchstone')
        # mplt.plot(np.transpose(freqs2), np.transpose(np.angle(Z11)), label='curve_fit')
        # mplt.xscale('log')
        # mplt.ylabel('Impedance Phase (degrees)')
        # mplt.xlabel('Frequency (log|Hz|)')
        # mplt.legend()
        
    def plot_residue_impedance(self, r, p, k, fstop=10e6):
        freqs = self.network.f
        s = 2j*3.14*freqs
        # assert len(r) == len(p), "\nPole and residue numbers dont match\n\n"
        # z11 = k
        # for i in range(len(r)):     
            # z11 = z11 + (r[i])/np.add(s,(-1)*p[i])
        y, x = signal.invres(r, p, k)
        # num = [y[0]*(s**4)+y[1]*(s**3)+y[2]*(s**2)+y[3]*s+y[4]]
        # den = [x[0]*(s**4)+x[1]*(s**3)+x[2]*(s**2)+x[3]*s+x[4]]
        num = 0
        den = 0
        print(f"len(y): {len(y)}")
        for i in range(len(y)):
            num += y[i]*(s**(len(y)-i-1))
        for j in range(len(x)):
            den += x[j]*(s**(len(x)-j-1))

        z11 = np.divide(num, den)
    
        freqs2 = np.reshape(freqs, (1,200))

        mplt.figure()
        mplt.plot(self.network.f, np.abs(self.network.z[:,0,0]), label='touchstone')
        mplt.plot(np.transpose(freqs2), np.transpose(np.abs(z11)), label='curve_fit')
        mplt.yscale('log')
        mplt.xscale('log')
        mplt.ylabel('Impedance Magnitude (log|Ohms|)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()
        
        mplt.figure()
        mplt.plot(self.network.f, np.angle(self.network.z[:,0,0], deg=True), label='touchstone')
        mplt.plot(np.transpose(freqs2), np.transpose(np.angle(z11, deg=True)), label='curve_fit')
        mplt.xscale('log')
        mplt.ylabel('Impedance Phase Error (degrees)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.legend()

        index_f10_mega_hz = np.where(self.network.f >= fstop)[0][0]
        db_error = 20*np.log10(np.transpose(np.abs(z11[:index_f10_mega_hz])))-20*np.log10(np.abs(self.network.z[:index_f10_mega_hz,0,0]))
        mplt.figure()
        mplt.plot(self.network.f[:index_f10_mega_hz], db_error, label='magnitude_error')
        # mplt.plot(np.transpose(freqs2), np.transpose(np.abs(z11)), label='curve_fit')
        # mplt.yscale('log')
        mplt.xscale('log')
        mplt.ylabel('Impedance Error (dB|Ohms|)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.grid(visible=True, which='both', axis='both')
        mplt.legend()
        
        phase_error = np.transpose(np.angle(z11[:index_f10_mega_hz], deg=True))-np.angle(self.network.z[:index_f10_mega_hz,0,0],deg=True)
        mplt.figure()
        mplt.plot(self.network.f[:index_f10_mega_hz], phase_error, label='phase_error')
        # mplt.plot(np.transpose(freqs2), np.transpose(np.angle(z11, deg=True)), label='curve_fit')
        mplt.xscale('log')
        mplt.ylabel('Impedance Phase Error (degrees)')
        mplt.xlabel('Frequency (log|Hz|)')
        mplt.grid(visible=True, which='both', axis='both')
        mplt.legend()

    def adice_rpc_to_impedance(a, b, k, z0, model_type=1):
        '''
        Converts adiceRPC (root polynomial coefficient) file parameters, (a,b,k) to an impedance Z(s).
        RPC coefficients are of the form M(s) = k + n_Σ {a[n] / (1 - s/b[n])}
                                      or M(s) = k + n_Σ {r[n] / (s - p[n])}, where n is the degree.
        M(s) is generic model parameter, which can either be S(s) or Y(s) [S11 or admittance (Y11)].
        Z(s) can be returned in residue or rational form.
        Residue  Z(s) = k + n_Σ {r[n] / (s - p[n])}
        Rational Z(s) = y(s)/x(s) where y[0] corresponds to s**n.

        Parameters
        -----------
        model_type: optional
            0 for Y parameter RPC coefficients, 1 for S parameter RPC coefficients.
        a
            List a[i] is equal to r[i]/p[i] (residues over poles) from the above M(s) equation.
        b
            List b[i] is equal to p[i] (poles) from M(s) equation.
        k
            DC term of M(s) equation.
        z0
            Characteristic impedance of S or Y parameter from RPC file.
        residue: optional
            1 to return residue version of Z, 0 to return rational version of Z

        Returns
        -----------
        y
            List of T(s) transfer function numerator coefficients.
        x
            List of T(s) transfer function denominator coefficients.  
        '''
        r = np.multiply(a, b)*(-1)
        y, x = signal.invres(r=r, p=b, k=k)
        s = sympy.symbols('s')  #Make variable s for j*omega
        m11_numerator = 0
        m11_denominator = 0
        for i in range(len(y)):
            m11_numerator += y[i]*(s**(len(y)-i-1))
        for j in range(len(x)):
            m11_denominator += x[j]*(s**(len(x)-j-1))
        if model_type == 1:     #If model type 1, S11->Z11.
            s11 = m11_numerator/m11_denominator
            z11_numerator = sympy.simplify(z0*(1+s11))
            z11_denominator = sympy.simplify(1-s11)
            z11 = z11_numerator/z11_denominator
        elif model_type == 0:   #If model type 0, Y11->Z11
            print(m11_denominator)
            z11 = m11_denominator/m11_numerator            
        else:
            print("Invalid model_type, please use 0 or 1.")
        z_num, z_den = sympy.fraction(z11)
        z_num_poly = sympy.Poly(z_num)
        z_den_poly = sympy.Poly(z_den)
        z_num_coeff = z_num_poly.all_coeffs()
        z_den_coeff = z_den_poly.all_coeffs()
        r, p, k = signal.residue(z_num_coeff, z_den_coeff)
        return z_num_coeff, z_num_coeff, r, p, k

    def adice_rpc_to_rational(a, b, k):
        '''
        Converts adiceRPC (root polynomial coefficient) file parameters, (a,b,k) to rational
        coefficients of the form T(s) = y(s)/x(s), whwere y[0] corresponds to s**n and n=degree.
        
        RPC coefficients are of the form M(s) = k + n_Σ {a[n] / (1 - s/b[n])}
                                      or M(s) = k + l_Σ {r[l] / (s - p[l])}
        M(s) is model type S(s) or Y(s) [S11 or admittance]
        
        Parameters
        -----------
        a
            List a[i] is equal to r[i]/p[i] (residues over poles) from the above M(s) equation.
        b
            List b[i] is equal to p[i] (poles) from M(s) equation.
        k
            DC term of M(s) equation.
        Returns
        -----------
        y
            List of T(s) transfer function numerator coefficients.
        x
            List of T(s) transfer function denominator coefficients.  
        '''
        r = np.multiply(a, b)*(-1)
        y, x = signal.invres(r=r, p=b, k=k)
        return y, x

        
    def dev_parse_RPC(input_file_name, port_num=1):
        pass
        #open file
        #file.readlines()
        #find modeltype
        #while first character is '.' keep iteratiting
    
    def dev_make_N_port_touchstone(self, output_file_name, num_ports):
        freq_obj = self.network.frequency
        freqs = self.network.f
        def make_port(name):
            return skrf.Circuit.Port(frequency=freq_obj, name=name,  z0=50)
        def make_resistor(name):
            return skrf.Circuit.SeriesImpedance(frequency=freq_obj, name=name,  z0=50, Z=50)
        connections = []
        resistors = []
        middle_resistor_node = []
        for i in range(num_ports): 
            resistors.append(make_resistor(f"resistor{i+1}"))
            connections.append([(make_port(f"port{i+1}"), 0), (resistors[i], 0)])
        for j in range(num_ports):
            middle_resistor_node.append((resistors[j], 1))
        connections.append(middle_resistor_node)
        # print(connections)
        # return
        circuit = skrf.Circuit(connections)
        network = circuit.network
        circuit.plot_graph(network_labels=True,
                           port_labels=True,
                           edge_labels=True,
                          )
        network.write_touchstone(file_name=output_file_name)
        return network
        # network_z11 = self.network.z[:,0,0]
        # ladder_z11 = network.z[:,0,0]
        # mplt.figure()
        # mplt.plot(freqs, np.abs(network_z11), label='touchstone')
        # mplt.plot(freqs, np.abs(ladder_z11),  label='1stage_ladder')
        # mplt.yscale('log')
        # mplt.xscale('log')
        # mplt.ylabel('Impedance Magnitude (log|Ohms|)')
        # mplt.xlabel('Frequency (log|Hz|)')
        # mplt.legend()
        

    ### Below are some functions that could get implemented ###
    
    def plot_Sbode(self):
        if self.network == None:
            print("No 2 port network stored, please input or create a network")
            return
        network = self.network
        network.plot_s_db(m=1, n=0, logx=True, z0=.1)

    def plot_Zbode(self):
        network = self.network
        if self.network == None:
            print("No 2 port network stored, please input or create a network")
            return
        network.plot_z_mag(m=0, n=0, logx=True)
        print(skrf.network.s2z(network.s, z0=.1))

    def model_fit_RC_LPF(self, touchstone_file_name):
        # turn touchstone to Bode plot
        # curve fit equation to get R and C, real version will have more complicated math with LL and RR ratios for ladder and constraints like @DC R1||R2||..||Rn = Rdc
        input_network = skrf.Network(touchstone_file_name)
        z0 = input_network.z0[0][0] #per freq,port array
        freqs = input_network.frequency
        
    def LR_modelfit(self, f_max=30e6): #This is Dave's algorithm
        model_params = {}
        model_z_thru = 1./self.network.y[:,0,0]
        model_params['Rdc'] = model_z_thru.real[0] #Lowest frequency
        i_max = np.where(self.network.frequency.f >= f_max)[0][0] #Find first index of f array that is greater than fmax
        r_tot_max = model_z_thru.real[i_max]
        model_params['R_skin_effect'] = (r_tot_max - model_params['Rdc'])/f_max**0.5 #2-point at min/max
        model_params['L'] = (model_z_thru.imag/self.network.frequency.w)[i_max//2] #middle freq fit
        print(model_params)
        
        
    def plot_smith_chart(self):
        '''
        Plots a smith chart to investigate the complex impedance over frequency.
        '''
        self.network.plot_s_smith()
        
    def get_resistor_skin_effect_model():
        '''
        Returns the approximate DC resistance and the scaling coefficient for the sqrt(2*pi*freq) skin effect term.
        '''
        return dc_res, ac_res_coefficient
        
    def get_series_LR(self):
        y_admittance = self.network.y[:,0,0]
        y_mag = abs(y_admittance)
        y_dc = y_mag[0]
        res = 1/y_dc
        print(res)


# if __name__ == '__main__':
    # print("Run delete_me.py")
    # customer_touchstone = 'C:\\Users\\BLeveret\\projects\\stowe_eval\\tests\\applications\\LT3390_4\\s_parameters_aptiv\\1V2.s4p'
    # tr = touchstone_utils(customer_touchstone)
    # mplt.show()